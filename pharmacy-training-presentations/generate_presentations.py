#!/usr/bin/env python3
"""Generate bilingual pharmacy team training PowerPoint presentations."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).parent / "presentations"
OUT.mkdir(exist_ok=True)

TEAL = RGBColor(0, 102, 102)
DARK = RGBColor(33, 37, 41)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(108, 117, 125)


def set_run(run, size=20, bold=False, color=DARK, font="Arial"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_title_slide(prs, title_en, title_ar, subtitle_en="", subtitle_ar=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TEAL
    bg.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title_en
    set_run(r, 36, True, WHITE)

    box2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(8.8), Inches(1.0))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = title_ar
    set_run(r2, 28, True, WHITE)

    if subtitle_en or subtitle_ar:
        box3 = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(8.8), Inches(1.0))
        tf3 = box3.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = f"{subtitle_en}\n{subtitle_ar}" if subtitle_ar else subtitle_en
        set_run(r3, 16, False, WHITE)


def add_section_slide(prs, title_en, title_ar):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, 0, Inches(2.8), prs.slide_width, Inches(1.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(1.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title_en
    set_run(r, 32, True, WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = title_ar
    set_run(r2, 24, True, WHITE)


def add_bullet_slide(prs, title_en, title_ar, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL
    header.line.fill.background()

    ht = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.9))
    htf = ht.text_frame
    hp = htf.paragraphs[0]
    hr = hp.add_run()
    hr.text = title_en
    set_run(hr, 24, True, WHITE)
    hp2 = htf.add_paragraph()
    hr2 = hp2.add_run()
    hr2.text = title_ar
    set_run(hr2, 18, True, WHITE)

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.5))
    btf = body.text_frame
    btf.word_wrap = True

    for i, (en, ar) in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.space_after = Pt(10)
        p.level = 0
        r = p.add_run()
        r.text = f"• {en}"
        set_run(r, 17, False, DARK)
        if ar:
            p2 = btf.add_paragraph()
            p2.space_after = Pt(14)
            p2.level = 1
            r2 = p2.add_run()
            r2.text = ar
            set_run(r2, 15, False, GRAY)


def add_closing_slide(prs, text_en, text_ar):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TEAL
    bg.line.fill.background()
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text_en
    set_run(r, 32, True, WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = text_ar
    set_run(r2, 26, True, WHITE)


def build_deck(filename, title_en, title_ar, subtitle_en, subtitle_ar, sections):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, title_en, title_ar, subtitle_en, subtitle_ar)

    for section_title_en, section_title_ar, slides in sections:
        add_section_slide(prs, section_title_en, section_title_ar)
        for slide_title_en, slide_title_ar, bullets in slides:
            add_bullet_slide(prs, slide_title_en, slide_title_ar, bullets)

    add_closing_slide(prs, "Questions & Discussion", "أسئلة ومناقشة")
    path = OUT / filename
    prs.save(str(path))
    print(f"Created: {path}")
    return path


from deck_content_part1 import ALL_DECKS as DECKS_1
from deck_content_part2 import ALL_DECKS_PART2 as DECKS_2
from deck_content_part3 import ALL_DECKS_PART3 as DECKS_3

ALL_DECKS = DECKS_1 + DECKS_2 + DECKS_3


def main():
    for deck in ALL_DECKS:
        build_deck(
            deck["filename"],
            deck["title_en"],
            deck["title_ar"],
            deck["subtitle_en"],
            deck["subtitle_ar"],
            deck["sections"],
        )
    print(f"\nDone — {len(ALL_DECKS)} presentations created in {OUT}")


if __name__ == "__main__":
    main()
